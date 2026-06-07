#!/usr/bin/env python3
"""
generate_ppt.py - Generate PowerPoint presentations from JSON spec using python-pptx.

Usage:
    python3 generate_ppt.py --spec slides.json --output presentation.pptx
    python3 generate_ppt.py --prompt "5-slide PPT about AI trends" --output ai.pptx
"""

import argparse
import sys
import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# ── Theme colors ───────────────────────────────────────
THEMES = {
    "blue": {
        "primary": RGBColor(0x1F, 0x45, 0x7C),
        "accent": RGBColor(0x44, 0x8A, 0xFF),
        "text": RGBColor(0x22, 0x22, 0x22),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xF5, 0xF7, 0xFA),
    },
    "green": {
        "primary": RGBColor(0x1B, 0x5E, 0x20),
        "accent": RGBColor(0x4C, 0xAF, 0x50),
        "text": RGBColor(0x22, 0x22, 0x22),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xF0, 0xF9, 0xF0),
    },
    "red": {
        "primary": RGBColor(0x8B, 0x00, 0x00),
        "accent": RGBColor(0xE8, 0x7A, 0x7A),
        "text": RGBColor(0x22, 0x22, 0x22),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xFF, 0xF5, 0xF5),
    },
    "dark": {
        "primary": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x7C, 0x3A, 0xED),
        "text": RGBColor(0xEE, 0xEE, 0xEE),
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "light_bg": RGBColor(0x2A, 0x2A, 0x4A),
    },
    "minimal": {
        "primary": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x88, 0x88, 0x88),
        "text": RGBColor(0x33, 0x33, 0x33),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "light_bg": RGBColor(0xFA, 0xFA, 0xFA),
    },
}


def get_theme(name):
    return THEMES.get(name, THEMES["blue"])


def set_cell_bg(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def format_text_frame(text_frame, font_size=18, bold=False, color=None):
    for para in text_frame.paragraphs:
        para.font.size = Pt(font_size)
        para.font.bold = bold
        if color:
            para.font.color.rgb = color


def add_title_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Top color bar
    top_bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        0, 0, prs.slide_width, Inches(1.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = theme["primary"]
    top_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8),
        prs.slide_width - Inches(1), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = spec.get("title", "Presentation Title")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    # Subtitle
    if spec.get("subtitle"):
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5),
            prs.slide_width - Inches(1), Inches(1)
        )
        sf = sub_box.text_frame
        p = sf.paragraphs[0]
        p.text = spec["subtitle"]
        p.font.size = Pt(20)
        p.font.color.rgb = theme["text"]

    # Author
    if spec.get("author"):
        auth_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8),
            prs.slide_width - Inches(1), Inches(0.5)
        )
        af = auth_box.text_frame
        p = af.paragraphs[0]
        p.text = f"By {spec['author']}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = theme["accent"]


def add_content_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec.get("title", "Content Slide")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    # Bullets
    bullets = spec.get("bullets", [])
    if bullets:
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3),
            prs.slide_width - Inches(1.6),
            prs.slide_height - Inches(2)
        )
        cf = content_box.text_frame
        cf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                para = cf.paragraphs[0]
            else:
                para = cf.add_paragraph()
            para.text = str(bullet)
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = theme["text"]
            para.space_after = Pt(8)

    # Free content string
    content = spec.get("content")
    if content and not bullets:
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3),
            prs.slide_width - Inches(1.6),
            prs.slide_height - Inches(2)
        )
        cf = content_box.text_frame
        cf.word_wrap = True
        p = cf.paragraphs[0]
        p.text = content
        p.font.size = Pt(16)
        p.font.color.rgb = theme["text"]


def add_two_column_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec.get("title", "Two Column Slide")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    left_items = spec.get("left", [])
    right_items = spec.get("right", [])
    col_width = (prs.slide_width - Inches(2)) / 2

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.3),
        col_width - Inches(0.2),
        prs.slide_height - Inches(2)
    )
    lf = left_box.text_frame
    lf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            para = lf.paragraphs[0]
        else:
            para = lf.add_paragraph()
        para.text = f"\u2022 {item}"
        para.font.size = Pt(16)
        para.font.color.rgb = theme["text"]
        para.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(0.8) + col_width + Inches(0.4), Inches(1.3),
        col_width - Inches(0.2),
        prs.slide_height - Inches(2)
    )
    rf = right_box.text_frame
    rf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            para = rf.paragraphs[0]
        else:
            para = rf.add_paragraph()
        para.text = f"\u2022 {item}"
        para.font.size = Pt(16)
        para.font.color.rgb = theme["text"]
        para.space_after = Pt(6)


def add_image_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec.get("title", "Image Slide")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    # Image
    img_path = spec.get("image")
    if img_path and os.path.exists(img_path):
        try:
            from PIL import Image
            with Image.open(img_path) as img:
                img_width, img_height = img.size
            aspect = img_width / img_height

            max_w = prs.slide_width - Inches(1)
            max_h = prs.slide_height - Inches(2)
            h = max_h
            w = h * aspect
            if w > max_w:
                w = max_w
                h = w / aspect

            left = (prs.slide_width - w) / 2
            top = Inches(1.2)
            slide.shapes.add_picture(img_path, left, top, width=w)
        except Exception:
            pass

    # Caption
    caption = spec.get("caption")
    if caption:
        cap_box = slide.shapes.add_textbox(
            Inches(0.5), prs.slide_height - Inches(0.8),
            prs.slide_width - Inches(1), Inches(0.5)
        )
        cf = cap_box.text_frame
        p = cf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = theme["accent"]


def add_table_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers or not rows:
        return

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        prs.slide_width - Inches(1), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = spec.get("title", "Table Slide")
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    nrows = len(rows) + 1
    ncols = len(headers)
    table_height = Inches(0.5) * nrows
    table_width = prs.slide_width - Inches(2)

    left = Inches(0.5)
    top = Inches(1.3)
    table = slide.shapes.add_table(nrows, ncols, left, top, table_width, table_height).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        set_cell_bg(cell, theme["primary"])
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.color.rgb = theme["text"]
            if r % 2 == 0:
                set_cell_bg(cell, theme["light_bg"])


def add_blank_slide(prs, spec, theme):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)


SLIDE_LAYOUTS = {
    "title": add_title_slide,
    "content": add_content_slide,
    "two-column": add_two_column_slide,
    "image": add_image_slide,
    "table": add_table_slide,
    "blank": add_blank_slide,
}


def generate_from_spec(spec, output_path, theme_name="blue"):
    """Generate PPT from JSON spec."""
    prs = Presentation()
    theme = get_theme(theme_name)

    spec_theme = spec.get("theme", theme_name)
    theme = get_theme(spec_theme)

    slides_spec = spec.get("slides", [])
    if not slides_spec:
        raise ValueError("No slides defined in spec")

    for slide_spec in slides_spec:
        layout = slide_spec.get("layout", "content")
        add_fn = SLIDE_LAYOUTS.get(layout, add_content_slide)
        add_fn(prs, slide_spec, theme)

    prs.save(output_path)
    return len(slides_spec)


def generate_from_prompt(prompt, output_path, theme_name="blue", slide_count=5):
    """Generate a basic PPT structure from a text prompt."""
    words = prompt.split()
    title = " ".join(words[:8]) if len(words) > 8 else prompt

    slides = [
        {
            "layout": "title",
            "title": title,
            "subtitle": "Generated by OpenJarvis AI",
            "author": "OpenJarvis",
        },
        {
            "layout": "content",
            "title": "Overview",
            "bullets": [
                f"Topic: {prompt}",
                "Key points will be expanded here",
                "Data and analysis included as needed",
            ],
        },
    ]

    for i in range(slide_count - 2):
        slides.append({
            "layout": "content",
            "title": f"Section {i+1}",
            "bullets": [
                f"Content for section {i+1}",
                "Add your details here",
                "Visual aids recommended",
            ],
        })

    slides.append({
        "layout": "content",
        "title": "Thank You",
        "bullets": ["Questions?", "Contact: user@example.com"],
    })

    spec = {"title": title, "theme": theme_name, "slides": slides}
    return generate_from_spec(spec, output_path, theme_name)


def main():
    parser = argparse.ArgumentParser(description="Generate PowerPoint from JSON spec or text prompt")
    parser.add_argument("--spec", default=None, help="Path to JSON spec file")
    parser.add_argument("--prompt", default=None, help="Text prompt to generate PPT from")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument("--theme", default="blue", choices=list(THEMES.keys()), help="Color theme")
    parser.add_argument("--slides", type=int, default=5, help="Number of slides when using --prompt")

    args = parser.parse_args()

    try:
        if args.spec:
            spec_path = Path(args.spec)
            if not spec_path.exists():
                print(f"Error: spec file not found: {args.spec}", file=sys.stderr)
                sys.exit(1)
            spec = json.loads(spec_path.read_text(encoding="utf-8"))
            count = generate_from_spec(spec, args.output, args.theme)
            print(f"Generated {count}-slide PPT: {args.output}")

        elif args.prompt:
            count = generate_from_prompt(args.prompt, args.output, args.theme, args.slides)
            print(f"Generated {count}-slide PPT from prompt: {args.output}")

        else:
            print("Error: Must provide either --spec or --prompt", file=sys.stderr)
            parser.print_help()
            sys.exit(1)

    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
